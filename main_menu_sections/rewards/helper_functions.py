
from datetime import datetime
import logging
import os
import re

from pptx import Presentation

from template_maker.file_exports import convert_ppt_to_image
from utils.database import execute_query, get_data


logging.basicConfig(
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s", level=logging.WARNING
)
logger = logging.getLogger(__name__)

def create_daily_gifts_folders(daily_gifts_path):
    """Creates the daily gifts folder structure (days 1-31)."""
    os.makedirs(daily_gifts_path)
    for day in range(1, 32):
        day_folder_path = os.path.join(daily_gifts_path, f"day_{day}")
        os.makedirs(day_folder_path)
        # Create a default message file for each day
        with open(
            os.path.join(day_folder_path, "message.txt"), "w", encoding="utf-8"
        ) as f:
            f.write(f"هذه هي هدية اليوم {day}! 🎁")


async def has_user_claimed_daily_reward(user_id):
    """Checks if the user has already claimed the daily reward today."""
    user_data = get_data(
        "SELECT last_daily_reward_claim FROM users WHERE telegram_id = ?", (user_id,)
    )

    if not user_data or not user_data[0][0]:  # No claim date or first time claim
        return False

    last_claim_date = datetime.strptime(user_data[0][0], "%Y-%m-%d").date()
    today = datetime.now().date()

    return last_claim_date == today


async def increment_user_daily_gifts_used(user_id):
    """Increments the user's number_of_daily_gifts_used in the database."""
    is_today = await has_user_claimed_daily_reward(user_id)
    if is_today:
        return
    execute_query(
        "UPDATE users SET number_of_daily_gifts_used = number_of_daily_gifts_used + 1, last_daily_reward_claim = ? WHERE telegram_id = ?",
        (datetime.now().strftime("%Y-%m-%d"), user_id),
    )

async def process_ppt_design_user_data(ppt_file, user_data):
    try:
        prs = Presentation(ppt_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            placeholders = re.findall(r"\((.*?)\)", text)
                            for placeholder in placeholders:
                                if placeholder in user_data:
                                    replacement_text = user_data[placeholder]

                                    # Preserve formatting: Replace run text directly
                                    run.text = text.replace(f"({placeholder})", replacement_text)

        modified_ppt = f"temp/temp_{user_data.get('telegram_id', 'user')}.pptx"
        prs.save(modified_ppt)

        image_path = f"temp/temp_{user_data.get('telegram_id', 'user')}.png"
        await convert_ppt_to_image(modified_ppt, image_path)

        os.remove(modified_ppt)
        return image_path
    except Exception as e:
        logger.error(f"Error processing PowerPoint design: {e}")
        raise

def get_user_custom_data(user_id):
    query = """
        SELECT name, phone_number, email, custom_text
        FROM user_customizations
        WHERE telegram_id = ?
    """
    result = get_data(query, (user_id,))
    if result:
        name, phone_number, email, custom_text = result[0]
        return {"telegram_id": user_id, "name": name, "phone_number": phone_number, "email": email, "custom_text": custom_text}
    else:
        query = "SELECT name FROM users WHERE telegram_id = ?"
        result = get_data(query, (user_id,))
        if result:
            name = result[0][0]
            return {"telegram_id": user_id, "name": name, "phone_number": "", "email": "", "custom_text": ""}
        else:
            return {"telegram_id": user_id}

async def get_user_stats(user_id):
    """Fetches and formats user statistics from the database."""
    try:
        user_data = get_data(
            "SELECT points, percentage_expected, usage_time, total_number_of_created_questions FROM users WHERE telegram_id = ?",
            (user_id,),
        )
    except Exception as e:
        logger.error(f"Error fetching user stats: {e}")
        return None

    if not user_data:
        return None

    points, percentage_expected, usage_time_str, questions_created = user_data[0]
    usage_time = convert_usage_time_to_hours(usage_time_str)

    return {
        "points": points if points is not None else 0,
        "percentage": (
            round(percentage_expected) if percentage_expected is not None else 0
        ),
        "time_spent": usage_time,
        "questions_created": questions_created if questions_created is not None else 0,
    }


def convert_usage_time_to_hours(usage_time_str):
    """Converts usage time string (HH:MM:SS) to hours (float)."""
    if usage_time_str:
        try:
            hours, minutes, seconds = map(int, usage_time_str.split(":"))
            return round(hours + minutes / 60 + seconds / 3600, 2)
        except ValueError:
            print(f"Invalid usage time format: {usage_time_str}")
            return 0
    return 0

def format_reward_message(user_value, target_value, reward_text, metric_name, unit):
    """Formats a single reward message."""
    if user_value >= target_value:
        return f"✅ - {reward_text}: حققت {user_value} {unit} من {target_value} {unit} المطلوبة من {metric_name}.\n"
    else:
        return f"⚠️ - {reward_text}: تحتاج إلى {target_value - user_value} {unit} إضافية من {metric_name} للوصول إلى {target_value} {unit}.\n"


def format_stats_text(user_stats, reward_messages: list):
    """Formats the statistics and reward messages into a single string."""
    reward_messages = "".join(reward_messages)

    return (
        f"📊 إحصائياتك:\n"
        f"🏅 نقاطك: {user_stats['points']}\n"
        f"📈 النسبة المئوية: {user_stats['percentage']}%\n"
        f"⏳ وقت الدراسة: {user_stats['time_spent']} ساعة\n"
        f"✏️ عدد الأسئلة التي قمت بإنشائها: {user_stats['questions_created']}\n\n"
        f"🎁 مكافآتك:\n{reward_messages}"
    )


def format_no_rewards_text(user_stats):
    """Formats the message when the user has no rewards."""
    return (
        f"📊 إحصائياتك:\n"
        f"🏅 نقاطك: {user_stats['points']}\n"
        f"📈 النسبة المئوية: {user_stats['percentage']}%\n"
        f"⏳ وقت الدراسة: {user_stats['time_spent']} ساعة\n"
        f"✏️ عدد الأسئلة التي قمت بإنشائها: {user_stats['questions_created']}\n\n"
        f"🙅‍♂️ لم تكسب أي مكافآت حتى الآن. استمر في الدراسة!"
    )


# --- Validation functions with character limits and user-friendly messages ---
def validate_name(name):
    if not name:
        return "الاسم لا يمكن أن يكون فارغًا."
    if len(name) < 2:
        return "الرجاء إدخال اسم صحيح (على الأقل حرفين)."
    if len(name) > 25:
        return "الاسم طويل جدًا. يجب أن يكون أقل من 25 حرفًا."
    return None

def validate_phone(phone):
    if not phone:
        return "رقم الهاتف لا يمكن أن يكون فارغًا."
    if not re.match(r"^\+?[1-9]\d{1,14}$", phone):
        return "الرجاء إدخال رقم هاتف صحيح (أرقام فقط، بحد أقصى 15 رقمًا)."
    return None

def validate_email(email):
    if not email:
        return "البريد الإلكتروني لا يمكن أن يكون فارغًا."
    if not re.match(r"[^@]+@[^@]+\.[^@]+", email):
        return "الرجاء إدخال بريد إلكتروني صحيح."
    if len(email) > 100:
       return "البريد الإلكتروني طويل جدًا. يجب أن يكون أقل من 100 حرفًا."
    return None

def validate_custom_text(text):
    if not text:
        return "النص لا يمكن أن يكون فارغًا."
    if len(text) > 200:
        return "النص طويل جدًا. يجب أن يكون أقل من 200 حرفًا."
    return None