import asyncio
from concurrent.futures import ThreadPoolExecutor
import copy
import logging
import platform
import shutil
from docx import Document
from pdf2image import convert_from_path
import subprocess
import os
import tempfile
import cv2
from docxtpl import DocxTemplate
from pptx import Presentation
import time
from docxcompose.composer import Composer
from docx.enum.section import WD_SECTION
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


async def generate_word_doc(template_path, output_path, quiz_data):
    """Generates the Word document."""
    try:
        doc = DocxTemplate(template_path)
        doc.render(quiz_data)
        doc.save(output_path)
    except Exception as e:
        logger.error(f"Error generating Word document: {e}")
        raise


async def merge_word_documents(main_doc_path, q_and_a_doc_path, output_path):
    """Merges the Main and Q\&A Word documents into a single document."""
    try:
        main_doc = Document(main_doc_path)
        # Add a page break to the end of the main document
        main_doc.add_section(WD_SECTION.NEW_PAGE)

        composer = Composer(main_doc)
        q_and_a_doc = Document(q_and_a_doc_path)
        composer.append(q_and_a_doc)

        composer.save(output_path)
    except Exception as e:
        logger.error(f"Error merging Word documents: {e}")
        raise


async def convert_docx_to_pdf(word_file, pdf_file=None):
    # Set default output name if pdf_file is not specified
    if pdf_file is None:
        pdf_file = os.path.splitext(word_file)[0] + ".pdf"

    # Set a temporary directory for the output to handle custom names
    temp_dir = os.path.dirname(pdf_file)
    temp_pdf_path = os.path.join(
        temp_dir, os.path.splitext(os.path.basename(word_file))[0] + ".pdf"
    )

    if platform.system() == "Windows":
        soffice_path = r"C:\Program Files\LibreOffice\program\soffice"
    else:
        soffice_path = "libreoffice"
    # Run the LibreOffice command to convert to PDF in temp directory
    try:
        proc = await asyncio.create_subprocess_exec(
            soffice_path,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            temp_dir,  # Use absolute path
            word_file,  # Use absolute path
        )
        await proc.communicate()
    except subprocess.CalledProcessError as e:
        logger.error(f"LibreOffice conversion failed: {e}")
        logger.error(f"Return code: {e.returncode}")
        logger.error(f"stderr: {e.stderr}")  # Print the error output from LibreOffice!
        return None
    except FileNotFoundError as e:
        logger.error(f"File not found during conversion: {e}")
        return None
    # Rename to the specified pdf_file name if it differs from temp_pdf_path
    if os.path.exists(temp_pdf_path):
        if temp_pdf_path != pdf_file:
            shutil.move(temp_pdf_path, pdf_file)
        print(f"Conversion successful: {pdf_file}")
        return pdf_file
    else:
        raise FileNotFoundError("PDF conversion failed.")

async def generate_powerpoint(main_template_path, q_and_a_template_path, output_path, quiz_data, user_data):
    """
    Generates a PowerPoint presentation, merges slides from two templates,
    and handles slide layouts and masters to avoid duplication.
    """
    try:
        prs = Presentation(main_template_path)
        q_and_a_prs = Presentation(q_and_a_template_path)

        # --- 1. Replace Placeholders in Main Template ---
        placeholder_mapping = {
            "studentName": user_data.get("studentName"),
            "phoneNumber": user_data.get("phoneNumber"),
            "expressionNumber": user_data.get("expressionNumber"),
            "modelNumber": user_data.get("modelNumber"),
            "date": user_data.get("date"),
            "questionsNumber": user_data.get("questionsNumber"),
            "studentsResults": user_data.get("studentsResults"),
        }

        replace_placeholders_in_slide(prs.slides, placeholder_mapping)

        # --- 2. Prepare Slide Layouts and Masters from Q&A Template ---
        q_and_a_layouts = {layout.name: layout for layout in q_and_a_prs.slide_layouts}
        q_and_a_masters = {master.name: master for master in q_and_a_prs.slide_masters}

        # --- 3. Add Q&A Slides ---
        questions = quiz_data["questions"]
        for i, question in enumerate(questions):
            for j in range(4):
                source_slide = q_and_a_prs.slides[j]

                # --- 3.1. Handle Slide Layout ---
                new_layout = get_or_create_slide_layout(prs, source_slide, q_and_a_layouts, q_and_a_masters)
                new_slide = prs.slides.add_slide(new_layout)

                # --- 3.2. Copy Shapes and Replace Placeholders---
                copy_shapes(source_slide, new_slide)
                copy_background(source_slide, new_slide)

                # Placeholder mapping for Q\&A slides
                question_placeholder_mapping = {
                    "QuestionNumber": str(question["QuestionNumber"]),
                    "MainCategoryName": question["MainCategoryName"],
                    "QuestionText": question["QuestionText"] if j == 0 else "",
                    "OptionA": question["OptionA"] if j == 2 else "",
                    "OptionB": question["OptionB"] if j == 2 else "",
                    "OptionC": question["OptionC"] if j == 2 else "",
                    "OptionD": question["OptionD"] if j == 2 else "",
                    "CorrectAnswer": question["CorrectAnswer"] if j == 3 else "",
                    "Explanation": question["Explanation"] if j == 3 else "",
                }

                replace_placeholders_in_slide([new_slide], question_placeholder_mapping)

                # --- 3.3 Copy Slide Transitions ---
                copy_transitions(source_slide, new_slide)

        prs.save(output_path)

    except Exception as e:
        print(f"Error in generate_powerpoint: {e}")
        raise

def replace_placeholders_in_slide(slides, placeholder_mapping):
    """Replaces placeholders in a slide with values from a dictionary."""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for placeholder, value in placeholder_mapping.items():
                        if placeholder == paragraph.text:
                            for run in paragraph.runs:
                                run.text = str(value)

def get_or_create_slide_layout(prs, source_slide, q_and_a_layouts, q_and_a_masters):
    """Gets an existing slide layout or creates a new one if it doesn't exist."""
    layout_name = source_slide.slide_layout.name
    if layout_name in q_and_a_layouts:
        # Check if this layout already exists in the main presentation
        existing_layout = None
        for existing_slide_layout in prs.slide_layouts:
            if existing_slide_layout.name == layout_name:
                existing_layout = existing_slide_layout
                break

        if not existing_layout:
            # If layout doesn't exist, add it from the Q\&A template
            source_master = source_slide.slide_layout.slide_master
            master_name = source_master.name

            if master_name not in [m.name for m in prs.slide_masters]:
                # If the master doesn't exist, copy it
                new_master = prs.slide_masters.add_master(source_master)
            else:
                # Master exists, find it
                for m in prs.slide_masters:
                    if m.name == master_name:
                        new_master = m
                        break

            # Add the layout to the master
            new_layout = new_master.slide_layouts.add_layout(q_and_a_layouts[layout_name])
        else:
            new_layout = existing_layout  # Use existing layout

        return new_layout
    else:
        print(f"Warning: Layout '{layout_name}' not found in Q\&A template.")
        return prs.slide_layouts[0]  # Fallback

def copy_shapes(source_slide, new_slide):
    """Copies shapes from the source slide to the new slide."""
    for shape in source_slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with open(shape.image.filename, 'rb') as image_file:
                    new_slide.shapes.add_picture(image_file, shape.left, shape.top, shape.width, shape.height)
            else:
                el = shape.element
                new_el = copy.deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        except Exception as e:
            print(f"Error copying shape: {e}")

def copy_background(source_slide, new_slide):
    """Copies the background from the source slide to the new slide."""
    try:
        new_slide.background.fill.solid()
        new_slide.background.fill.fore_color.rgb = source_slide.background.fill.fore_color.rgb
    except Exception as e:
        print(f"Error copying background: {e}")

def copy_transitions(source_slide, new_slide):
    """Copies transitions from the source slide to the new slide."""
    source_xml = source_slide._element
    new_xml = new_slide._element
    source_transition = source_xml.find(
        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}transition"
    )
    if source_transition is not None:
        existing_transition = new_xml.find(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}transition"
        )
        if existing_transition is not None:
            new_xml.remove(existing_transition)
        new_xml.append(copy.deepcopy(source_transition))


def _convert_with_windows(pptx_path, mp4_path):
    """
    Blocking conversion for Windows using COM automation.
    This function runs in a separate thread.
    """
    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()  # Initialize COM for the current thread
    try:
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        pptx_path = os.path.abspath(pptx_path)
        mp4_path = os.path.abspath(mp4_path)

        presentation = ppt.Presentations.Open(pptx_path, WithWindow=False)
        presentation.CreateVideo(mp4_path, -1, 4, 1080, 24, 60)
        start_time_stamp = time.time()

        while True:
            time.sleep(4)  # Blocking in the thread
            try:
                os.rename(mp4_path, mp4_path)
                print("Success")
                break
            except Exception as e:
                print(f"Waiting for video creation: {e}")

        end_time_stamp = time.time()
        print(f"Time taken: {end_time_stamp - start_time_stamp}")
        presentation.Close()
        ppt.Quit()
    finally:
        pythoncom.CoUninitialize()  # Clean up COM initialization

async def convert_pptx_to_mp4(
    pptx_path, mp4_path, fps=0.5, dpi=300, image_format="png"
):
    """
    Converts a PPTX file to an MP4 video by first converting it to a PDF,
    then converting each PDF page to an image, and finally combining
    the images into a video.

    Parameters:
    - pptx_path: Path to the input PPTX file
    - mp4_path: Path to the output MP4 file
    - fps: Frames per second for the output video
    - dpi: DPI setting for converting PDF to images
    - image_format: Format to save images (default is PNG)
    """
    if platform.system() == "Windows":
        loop = asyncio.get_event_loop()
        try:
            with ThreadPoolExecutor() as executor:
                # Use a thread pool for blocking calls
                await loop.run_in_executor(executor, _convert_with_windows, pptx_path, mp4_path)
        except Exception as e:
            print(f"Error in convert_pptx_to_mp4_with_windows: {e}")

    else:
        # Step 1: Convert PPTX to PDF using LibreOffice
        with tempfile.TemporaryDirectory() as temp_dir:
            pdf_filename = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
            pdf_path = os.path.join(temp_dir, pdf_filename)

            if platform.system() == "Windows":
                soffice_path = r"C:\Program Files\LibreOffice\program\soffice"
            else:
                soffice_path = "libreoffice"
            # Run the LibreOffice command to convert PPTX to PDF
            result = subprocess.run(
                [
                    soffice_path,
                    "--headless",
                    "--invisible",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    temp_dir,
                    pptx_path,
                ],
                check=True,
            )

            # Check if the PDF was created successfully
            if not os.path.exists(pdf_path):
                raise FileNotFoundError("PDF conversion failed.")
            print(f"Converted PPTX to PDF: {pdf_path}")

            # Step 2: Convert PDF to images
            if platform.system() == "Windows":
                # Step 2: Convert PDF to images asynchronously
                images = await asyncio.to_thread(
                    convert_from_path,
                    pdf_path,
                    dpi=dpi,
                    fmt=image_format,
                    poppler_path=r"C:\poppler-24.08.0\Library\bin",
                )
            else:
                images = await asyncio.to_thread(
                    convert_from_path, pdf_path, dpi=dpi, fmt=image_format
                )

            # Step 3: Save images to temporary files
            image_paths = []
            for i, img in enumerate(images):
                img_path = os.path.join(
                    temp_dir, f"page_{str(i + 1).zfill(3)}.{image_format}"
                )
                img.save(img_path, image_format.upper())
                image_paths.append(img_path)

            print(f"PDF converted to images: {len(image_paths)} images generated.")

            # Step 4: Create video from images
            # Load first image to get frame dimensions
            frame = cv2.imread(image_paths[0])
            height, width, layers = frame.shape

            # Define the video codec and create VideoWriter object
            fourcc = cv2.VideoWriter_fourcc(*"mp4v")  # Codec for mp4
            video = cv2.VideoWriter(mp4_path, fourcc, fps, (width, height))

            # Add each image to the video
            for img_path in image_paths:
                frame = cv2.imread(img_path)
                video.write(frame)  # Write each image as a frame

            # Release the video writer
            video.release()
            print(f"Video saved as {mp4_path}")

    # Temporary files (PDF and images) are automatically deleted with temp_dir


def convert_pptx_to_mp4_with_windows(
    pptx_path, mp4_path, fps=0.5, dpi=300, image_format="png"
):
    try:
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        pptx_path = os.path.abspath(pptx_path)
        mp4_path = os.path.abspath(mp4_path)
        presentation = ppt.Presentations.Open(pptx_path, WithWindow=False)
        presentation.CreateVideo(mp4_path, -1, 4, 1080, 24, 60)
        start_time_stamp = time.time()

        while True:
            time.sleep(4)
            try:
                os.rename(mp4_path, mp4_path)
                print("Success")
                break
            except Exception as e:
                print(f"Waiting for video creation: {e}")

        end_time_stamp = time.time()
        print(end_time_stamp - start_time_stamp)
        presentation.Close()
        ppt.Quit()

    except Exception as e:
        print(f"Error in convert_pptx_to_mp4_with_windows: {e}")


async def convert_ppt_to_image(ppt_file_path, img_path, dpi=300, image_format="png"):
    """
    Converts a PPTX file to an image.
    """
    try:
        # Step 1: Convert PPTX to PDF using LibreOffice asynchronously
        with tempfile.TemporaryDirectory() as temp_dir:
            pdf_filename = os.path.splitext(os.path.basename(ppt_file_path))[0] + ".pdf"
            pdf_path = os.path.join(temp_dir, pdf_filename)

            if platform.system() == "Windows":
                soffice_path = r"C:\Program Files\LibreOffice\program\soffice"
            else:
                soffice_path = "libreoffice"
            # Run the LibreOffice command asynchronously
            cmd = [
                soffice_path,
                "--headless",
                "--invisible",
                "--convert-to",
                "pdf",
                "--outdir",
                temp_dir,
                ppt_file_path,
            ]

            process = await asyncio.create_subprocess_exec(
                *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await process.communicate()

            if process.returncode != 0 or not os.path.exists(pdf_path):
                raise FileNotFoundError(f"PDF conversion failed: {stderr.decode()}")

            if platform.system() == "Windows":
                # Step 2: Convert PDF to images asynchronously
                images = await asyncio.to_thread(
                    convert_from_path,
                    pdf_path,
                    dpi=dpi,
                    fmt=image_format,
                    poppler_path=r"C:\poppler-24.08.0\Library\bin",
                )
            else:
                images = await asyncio.to_thread(
                    convert_from_path, pdf_path, dpi=dpi, fmt=image_format
                )

            # Save the first page as the image for simplicity
            images[0].save(img_path, image_format.upper())

            return img_path

    except Exception as e:
        logger.error(f"Error converting PPTX to image: {e}")
        raise
